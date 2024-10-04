from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_file, pptx_file):
    # Convertir PDF a imágenes
    images = convert_from_path(pdf_file)
    
    # Crear una nueva presentación de PowerPoint
    presentation = Presentation()
    
    for image in images:
        # Agregar una nueva diapositiva
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout vacío

        # Guardar la imagen temporalmente
        image_path = "temp_image.jpg"
        image.save(image_path, "JPEG")

        # Añadir la imagen a la diapositiva
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
        
        # Eliminar la imagen temporal
        os.remove(image_path)
    
    # Guardar la presentación como archivo PPTX
    presentation.save(pptx_file)
    print(f"Presentación guardada como {pptx_file}")

def main():
    # Pedir al usuario la ubicación del PDF
    pdf_file = input("Por favor, introduce la ruta completa del archivo PDF que deseas convertir: ")

    # Verificar que el archivo existe
    if not os.path.isfile(pdf_file):
        print("El archivo PDF no existe. Por favor, revisa la ruta e inténtalo de nuevo.")
        return

    # Pedir al usuario la ubicación de guardado para el archivo PPTX
    pptx_file = input("Por favor, introduce la ruta completa donde deseas guardar el archivo PPTX: ")

    # Añadir extensión .pptx si no está presente
    if not pptx_file.endswith(".pptx"):
        pptx_file += ".pptx"

    # Convertir PDF a PPTX
    pdf_to_pptx(pdf_file, pptx_file)

if __name__ == "__main__":
    main()
